"""
THERESE v2 - Code Executor pour Skills Office

Exécution sandboxée de code Python généré par le LLM.
Approche code-execution : LLM -> code Python -> exécution -> fichier.
"""

import ast
import asyncio
import logging
import re
from abc import abstractmethod
from pathlib import Path
from typing import Any

from app.services.skills.base import BaseSkill, SkillParams, SkillResult

logger = logging.getLogger(__name__)

# Timeout d'exécution en secondes
EXECUTION_TIMEOUT = 30

# Palette Synoptia pour injection dans le namespace
SYNOPTIA_COLORS = {
    "background": "#0B1226",
    "surface": "#131B35",
    "text": "#E6EDF7",
    "muted": "#A9B8D8",
    "primary": "#2451FF",
    "accent_cyan": "#22D3EE",
    "accent_magenta": "#E11D8D",
    "header_bg": "0F1E6D",
    "header_text": "E6EDF7",
    "row_alt": "F5F7FA",
    "input_blue": "3B82F6",
    "formula_black": "1A1A2E",
    "link_green": "22C55E",
    "heading": "0F1E6D",
    "body": "1A1A2E",
}

# Patterns bloqués dans le code généré
BLOCKED_PATTERNS: list[str] = [
    r"\bos\.",
    r"\bos\b\s*\(",
    r"\bsys\.",
    r"\bsys\b\s*\(",
    r"\bsubprocess\b",
    r"\bshutil\b",
    r"\bsocket\b",
    r"\brequests\b",
    r"\burllib\b",
    r"\b__import__\b",
    r"\beval\s*\(",
    r"\bexec\s*\(",
    r"\bcompile\s*\(",
    r"\bglobals\s*\(",
    r"\blocals\s*\(",
    r"\bgetattr\s*\(",
    r"\bsetattr\s*\(",
    r"\bdelattr\s*\(",
    r"\bbreakpoint\s*\(",
    r"\binput\s*\(",
]

# Imports autorisés par format
ALLOWED_IMPORTS: dict[str, set[str]] = {
    "xlsx": {
        "openpyxl",
        "openpyxl.Workbook",
        "openpyxl.styles",
        "openpyxl.styles.Font",
        "openpyxl.styles.PatternFill",
        "openpyxl.styles.Alignment",
        "openpyxl.styles.Border",
        "openpyxl.styles.Side",
        "openpyxl.styles.numbers",
        "openpyxl.chart",
        "openpyxl.chart.BarChart",
        "openpyxl.chart.LineChart",
        "openpyxl.chart.PieChart",
        "openpyxl.chart.Reference",
        "openpyxl.utils",
        "openpyxl.utils.get_column_letter",
        "pandas",
        "datetime",
        "json",
        "re",
        "math",
        "decimal",
        "decimal.Decimal",
        "time",
        "random",
        "copy",
        "string",
        "textwrap",
        "itertools",
        "collections",
    },
    "docx": {
        "docx",
        "docx.Document",
        "docx.shared",
        "docx.shared.Cm",
        "docx.shared.Pt",
        "docx.shared.Inches",
        "docx.shared.RGBColor",
        "docx.enum.text",
        "docx.enum.text.WD_ALIGN_PARAGRAPH",
        "docx.enum.style",
        "docx.enum.style.WD_STYLE_TYPE",
        "docx.enum.table",
        "docx.enum.table.WD_TABLE_ALIGNMENT",
        "docx.enum.section",
        "docx.oxml.ns",
        "docx.oxml.ns.qn",
        "datetime",
        "json",
        "re",
        "math",
        "decimal",
        "decimal.Decimal",
        "time",
        "random",
        "copy",
        "string",
        "textwrap",
        "itertools",
        "collections",
    },
    "pptx": {
        "pptx",
        "pptx.Presentation",
        "pptx.util",
        "pptx.util.Inches",
        "pptx.util.Pt",
        "pptx.util.Cm",
        "pptx.util.Emu",
        "pptx.dml.color",
        "pptx.dml.color.RGBColor",
        "pptx.enum.text",
        "pptx.enum.text.PP_ALIGN",
        "pptx.enum.text.MSO_ANCHOR",
        "pptx.enum.shapes",
        "pptx.enum.chart",
        "datetime",
        "json",
        "re",
        "math",
        "decimal",
        "decimal.Decimal",
        "time",
        "random",
        "copy",
        "string",
        "textwrap",
        "itertools",
        "collections",
    },
}


class CodeExecutionError(Exception):
    """Erreur lors de l'exécution du code généré."""

    pass


# Seuils minimum de contenu pour valider un document généré par code-execution.
MIN_CONTENT_ELEMENTS = {
    "docx": 3,   # au moins 3 paragraphes non vides (hors titre)
    "pptx": 2,   # au moins 2 slides
    "xlsx": 2,   # au moins 2 lignes de données (hors header)
}


def _validate_document_content(output_path: str, format_type: str) -> bool:
    """
    Vérifie qu'un document généré contient assez de contenu.

    Retourne True si le document est suffisamment riche, False sinon.
    """
    min_elements = MIN_CONTENT_ELEMENTS.get(format_type, 2)
    path = Path(output_path)

    try:
        if format_type == "docx":
            from docx import Document
            doc = Document(str(path))
            non_empty = sum(1 for p in doc.paragraphs if p.text.strip())
            logger.debug(
                "Validation DOCX : %d paragraphes non vides (min=%d)",
                non_empty, min_elements,
            )
            return non_empty >= min_elements

        elif format_type == "pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            slide_count = len(prs.slides)
            logger.debug(
                "Validation PPTX : %d slides (min=%d)",
                slide_count, min_elements,
            )
            return slide_count >= min_elements

        elif format_type == "xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(path), read_only=True)
            ws = wb.active
            if ws is None:
                return False
            row_count = sum(
                1 for row in ws.iter_rows(max_row=100) if any(c.value for c in row)
            )
            wb.close()
            logger.debug(
                "Validation XLSX : %d lignes avec données (min=%d)",
                row_count, min_elements,
            )
            return row_count >= min_elements

    except Exception as e:
        logger.warning("Erreur validation contenu %s : %s", format_type, e)

    return True  # En cas d'erreur de lecture, on accepte le fichier


def extract_python_code(llm_response: str) -> str | None:
    """
    Extrait le code Python d'un bloc ```python``` dans la réponse LLM.

    Gère aussi les réponses tronquées où le ``` fermant est absent
    (ex: réponses longues coupées par max_tokens).

    Args:
        llm_response: Réponse complète du LLM

    Returns:
        Code Python extrait ou None si pas de bloc trouvé
    """
    # 1. Chercher un bloc complet ```python ... ```
    pattern = r"```python\s*\n(.*?)```"
    match = re.search(pattern, llm_response, re.DOTALL)
    if match:
        code = match.group(1).strip()
        if code:
            return code

    # 2. Essayer aussi ```py ... ```
    pattern_py = r"```py\s*\n(.*?)```"
    match_py = re.search(pattern_py, llm_response, re.DOTALL)
    if match_py:
        code = match_py.group(1).strip()
        if code:
            return code

    # 3. Fallback : bloc ouvert sans ``` fermant (réponse tronquée)
    pattern_open = r"```python\s*\n(.*)"
    match_open = re.search(pattern_open, llm_response, re.DOTALL)
    if match_open:
        code = match_open.group(1).strip()
        if code:
            logger.warning(
                "Code Python extrait sans bloc fermant (réponse tronquée, %d chars)",
                len(code),
            )
            return code

    # 4. Idem pour ```py sans fermant
    pattern_py_open = r"```py\s*\n(.*)"
    match_py_open = re.search(pattern_py_open, llm_response, re.DOTALL)
    if match_py_open:
        code = match_py_open.group(1).strip()
        if code:
            logger.warning(
                "Code Python extrait sans bloc fermant (réponse tronquée, %d chars)",
                len(code),
            )
            return code

    # 5. Dernier recours : détecter du code Python brut (sans bloc markdown)
    # Certains modèles (ex: Gemini) renvoient du code sans ```python```
    # On cherche des marqueurs forts de code python-docx/pptx/openpyxl
    code_markers = [
        r'from\s+(?:docx|pptx|openpyxl)\s+import',
        r'(?:Document|Presentation|Workbook)\s*\(',
        r'\.add_slide\s*\(',
        r'\.add_heading\s*\(',
        r'\.save\s*\(',
    ]
    marker_count = sum(1 for m in code_markers if re.search(m, llm_response))
    if marker_count >= 2:
        # Extraire les lignes qui ressemblent à du code Python
        # (ignorer les lignes de texte libre avant/après le code)
        lines = llm_response.split('\n')
        code_lines: list[str] = []
        in_code = False
        for line in lines:
            stripped = line.strip()
            # Début du code : import ou assignation ou commentaire Python
            if not in_code:
                if (stripped.startswith(('from ', 'import ', '#'))
                    or re.match(r'^[a-zA-Z_]\w*\s*=\s*', stripped)):
                    in_code = True
                    code_lines.append(line)
            else:
                # Fin du code : ligne vide après du contenu, ou texte narratif
                if stripped and not stripped.startswith('#') and not any(c in stripped for c in '=()[]{}.:,+-*/_"\'\\') and len(stripped.split()) > 5:
                    # Ligne de texte narratif (pas de code) - on arrête si on a déjà du code
                    if len(code_lines) > 5:
                        break
                else:
                    code_lines.append(line)

        if code_lines:
            code = '\n'.join(code_lines).strip()
            try:
                ast.parse(code)
                logger.warning(
                    "Code Python détecté sans bloc markdown (%d lignes, %d markers)",
                    len(code_lines),
                    marker_count,
                )
                return code
            except SyntaxError:
                # Tenter réparation
                repaired = repair_truncated_code(code)
                if repaired:
                    logger.warning(
                        "Code Python sans bloc markdown réparé (%d lignes)",
                        len(code_lines),
                    )
                    return repaired

    return None


def repair_truncated_code(code: str) -> str | None:
    """
    Tente de réparer du code Python tronqué en retirant les lignes
    incomplètes à la fin jusqu'à obtenir un code syntaxiquement valide.

    Si la réparation supprime l'appel .save(output_path), celui-ci est
    rajouté automatiquement en détectant le nom de la variable du document
    (wb, doc, prs, document, workbook, presentation).

    Args:
        code: Code Python potentiellement tronqué

    Returns:
        Code réparé ou None si impossible
    """
    # D'abord vérifier si le code est déjà valide
    try:
        ast.parse(code)
        return _ensure_save_call(code)
    except SyntaxError:
        pass

    # Retirer les lignes une par une depuis la fin
    lines = code.split("\n")
    for i in range(len(lines), 0, -1):
        candidate = "\n".join(lines[:i]).rstrip()
        if not candidate:
            continue
        try:
            ast.parse(candidate)
            removed = len(lines) - i
            logger.info(
                "Code tronqué réparé : %d lignes retirées sur %d",
                removed,
                len(lines),
            )
            # Ajouter .save(output_path) si absent après réparation
            candidate = _ensure_save_call(candidate)
            return candidate
        except SyntaxError:
            continue

    return None


def _ensure_save_call(code: str) -> str:
    """
    Vérifie que le code contient un appel .save(output_path).
    Si absent, détecte le nom de la variable du document et l'ajoute.

    Args:
        code: Code Python syntaxiquement valide

    Returns:
        Code avec .save(output_path) garanti
    """
    # Vérifier si un .save() existe déjà (avec n'importe quel argument)
    if re.search(r'\.save\s*\(', code):
        return code

    # Détecter le nom de la variable du document principal
    # Patterns courants : wb = Workbook(), doc = Document(), prs = Presentation()
    # Inclut aussi load_workbook() pour les fichiers Excel existants
    doc_var = None
    for pattern in [
        r"(\w+)\s*=\s*Workbook\s*\(",
        r"(\w+)\s*=\s*load_workbook\s*\(",
        r"(\w+)\s*=\s*Document\s*\(",
        r"(\w+)\s*=\s*Presentation\s*\(",
    ]:
        match = re.search(pattern, code)
        if match:
            doc_var = match.group(1)
            break

    if doc_var:
        save_line = f"\n{doc_var}.save(output_path)\n"
        logger.warning(
            "Appel .save(output_path) manquant, "

            "ajout automatique : %s.save(output_path)",
            doc_var,
        )
        return code + save_line

    return code


def validate_code(code: str) -> tuple[bool, str]:
    """
    Valide la sécurité du code Python généré.

    Vérifie :
    1. Syntaxe Python valide (via ast.parse)
    2. Pas de patterns dangereux (os, sys, subprocess, etc.)
    3. Pas d'imports non autorisés

    Args:
        code: Code Python à valider

    Returns:
        Tuple (est_valide, message_erreur)
    """
    # 1. Vérifier la syntaxe
    try:
        ast.parse(code)
    except SyntaxError as e:
        return False, f"Erreur de syntaxe Python : {e}"

    # 2. Vérifier les patterns bloqués
    for pattern in BLOCKED_PATTERNS:
        if re.search(pattern, code):
            return False, f"Pattern interdit détecté : {pattern}"

    # 3. Vérifier que open() n'est utilisé qu'avec output_path
    # On autorise open() uniquement via les appels .save() des bibliothèques
    open_calls = re.findall(r"\bopen\s*\(", code)
    if open_calls:
        # Vérifier que open() est utilisé uniquement avec output_path
        # Pattern autorisé : open(output_path, ...) ou open(str(output_path), ...)
        safe_open = re.findall(
            r"\bopen\s*\(\s*(?:str\s*\(\s*)?output_path", code
        )
        if len(open_calls) != len(safe_open):
            return False, "open() n'est autorisé qu'avec output_path"

    return True, ""


def _validate_imports(code: str, format_type: str) -> tuple[bool, str]:
    """
    Vérifie que les imports sont autorisés pour le format donné.

    Args:
        code: Code Python à valider
        format_type: Format du fichier (xlsx, docx, pptx)

    Returns:
        Tuple (est_valide, message_erreur)
    """
    allowed = ALLOWED_IMPORTS.get(format_type, set())

    try:
        tree = ast.parse(code)
    except SyntaxError:
        return False, "Erreur de syntaxe"

    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                module_root = alias.name.split(".")[0]
                if module_root not in {m.split(".")[0] for m in allowed}:
                    return False, f"Import interdit : {alias.name}"
        elif isinstance(node, ast.ImportFrom) and node.module:
            module_root = node.module.split(".")[0]
            if module_root not in {m.split(".")[0] for m in allowed}:
                return False, f"Import interdit : from {node.module}"

    return True, ""


def _build_namespace(
    output_path: str, title: str, format_type: str, nb_slides: int = 10
) -> dict[str, Any]:
    """
    Construit le namespace d'exécution sandboxé.

    Args:
        output_path: Chemin de sauvegarde du fichier
        title: Titre du document
        format_type: Format du fichier (xlsx, docx, pptx)

    Returns:
        Namespace dict avec les variables et imports autorisés
    """
    import datetime
    import json
    import math
    import re as re_module
    from decimal import Decimal

    namespace: dict[str, Any] = {
        # Variables injectées
        "output_path": output_path,
        "title": title,
        "nb_slides": nb_slides,
        "SYNOPTIA_COLORS": SYNOPTIA_COLORS.copy(),
        # Modules communs
        "datetime": datetime,
        "json": json,
        "re": re_module,
        "math": math,
        "Decimal": Decimal,
        # Builtins restreints
        "__builtins__": {
            "print": print,
            "len": len,
            "range": range,
            "enumerate": enumerate,
            "zip": zip,
            "map": map,
            "filter": filter,
            "sorted": sorted,
            "reversed": reversed,
            "list": list,
            "dict": dict,
            "set": set,
            "tuple": tuple,
            "str": str,
            "int": int,
            "float": float,
            "bool": bool,
            "max": max,
            "min": min,
            "sum": sum,
            "abs": abs,
            "round": round,
            "isinstance": isinstance,
            "type": type,
            "hasattr": hasattr,
            "None": None,
            "True": True,
            "False": False,
            "ValueError": ValueError,
            "TypeError": TypeError,
            "KeyError": KeyError,
            "IndexError": IndexError,
            "Exception": Exception,
            "__import__": _restricted_import(format_type),
        },
    }

    # Imports spécifiques selon le format
    if format_type == "xlsx":
        import openpyxl
        from openpyxl import Workbook
        from openpyxl.chart import BarChart, LineChart, PieChart, Reference
        from openpyxl.styles import (
            Alignment,
            Border,
            Font,
            PatternFill,
            Side,
        )
        from openpyxl.utils import get_column_letter

        namespace.update(
            {
                "openpyxl": openpyxl,
                "Workbook": Workbook,
                "Font": Font,
                "PatternFill": PatternFill,
                "Alignment": Alignment,
                "Border": Border,
                "Side": Side,
                "BarChart": BarChart,
                "LineChart": LineChart,
                "PieChart": PieChart,
                "Reference": Reference,
                "get_column_letter": get_column_letter,
            }
        )
    elif format_type == "docx":
        import docx
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Cm, Inches, Pt, RGBColor

        namespace.update(
            {
                "docx": docx,
                "Document": Document,
                "Cm": Cm,
                "Pt": Pt,
                "Inches": Inches,
                "RGBColor": RGBColor,
                "WD_ALIGN_PARAGRAPH": WD_ALIGN_PARAGRAPH,
            }
        )
    elif format_type == "pptx":
        import pptx
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Cm, Emu, Inches, Pt

        namespace.update(
            {
                "pptx": pptx,
                "Presentation": Presentation,
                "Inches": Inches,
                "Pt": Pt,
                "Cm": Cm,
                "Emu": Emu,
                "RGBColor": RGBColor,
                "PP_ALIGN": PP_ALIGN,
                "MSO_ANCHOR": MSO_ANCHOR,
            }
        )

    return namespace


def _restricted_import(format_type: str):
    """
    Crée une fonction __import__ restreinte aux modules autorisés.

    Args:
        format_type: Format du fichier (xlsx, docx, pptx)

    Returns:
        Fonction __import__ sécurisée
    """
    import builtins

    allowed_roots = {m.split(".")[0] for m in ALLOWED_IMPORTS.get(format_type, set())}

    def safe_import(name, *args, **kwargs):
        root = name.split(".")[0]
        if root not in allowed_roots:
            raise ImportError(
                f"Import interdit : '{name}'. "
                f"Modules autorisés : {sorted(allowed_roots)}"
            )
        return builtins.__import__(name, *args, **kwargs)

    return safe_import


async def execute_sandboxed(
    code: str,
    output_path: str,
    title: str,
    format_type: str,
    nb_slides: int = 10,
) -> None:
    """
    Exécute du code Python dans un namespace restreint avec timeout.

    Args:
        code: Code Python à exécuter
        output_path: Chemin de sauvegarde du fichier
        title: Titre du document
        format_type: Format du fichier (xlsx, docx, pptx)

    Raises:
        CodeExecutionError: Si l'exécution échoue
    """
    # 1. Valider la sécurité du code
    is_valid, error_msg = validate_code(code)
    if not is_valid:
        raise CodeExecutionError(f"Validation échouée : {error_msg}")

    # 2. Valider les imports
    is_valid_imports, import_error = _validate_imports(code, format_type)
    if not is_valid_imports:
        raise CodeExecutionError(f"Validation imports échouée : {import_error}")

    # 3. Construire le namespace
    namespace = _build_namespace(output_path, title, format_type, nb_slides)

    # 4. Exécuter dans un thread avec timeout
    def _execute():
        try:
            compiled = compile(code, "<llm_generated>", "exec")
            exec(compiled, namespace)  # noqa: S102
        except Exception as e:
            raise CodeExecutionError(
                f"Erreur d'exécution : {type(e).__name__}: {e}"
            ) from e

    try:
        await asyncio.wait_for(
            asyncio.to_thread(_execute),
            timeout=EXECUTION_TIMEOUT,
        )
    except asyncio.TimeoutError:
        raise CodeExecutionError(
            f"Timeout : l'exécution a dépassé {EXECUTION_TIMEOUT}s"
        )
    except CodeExecutionError:
        raise
    except Exception as e:
        raise CodeExecutionError(
            f"Erreur inattendue : {type(e).__name__}: {e}"
        ) from e


class CodeGenSkill(BaseSkill):
    """
    Classe abstraite pour les skills qui génèrent du code Python.

    Étend BaseSkill avec une approche code-execution :
    1. Le LLM génère du code Python dans un bloc ```python```
    2. Le code est validé et exécuté dans une sandbox
    3. Si échec ou pas de code, fallback vers l'ancien parser
    """

    async def execute(self, params: SkillParams) -> SkillResult:
        """
        Exécute le skill avec approche code-execution.

        1. Extraire code Python de params.content
        2. Si code trouvé -> exécuter dans sandbox
        3. Si échec ou pas de code -> _fallback_execute()
        4. Vérifier que le fichier a été créé (taille > 0)

        Args:
            params: Paramètres de génération

        Returns:
            Résultat avec chemin vers le fichier généré
        """
        file_id = self.generate_file_id()
        output_path = self.get_output_path(file_id, params.title)

        # Tenter l'extraction du code Python
        code = extract_python_code(params.content)

        if code:
            # Si le code a une erreur de syntaxe, tenter une réparation
            try:
                ast.parse(code)
            except SyntaxError:
                logger.warning(
                    f"[{self.skill_id}] Code extrait avec erreur de syntaxe, "
                    f"tentative de réparation..."
                )
                repaired = repair_truncated_code(code)
                if repaired:
                    code = repaired
                else:
                    logger.warning(
                        f"[{self.skill_id}] Réparation impossible, "
                        f"fallback vers parser legacy"
                    )
                    code = None

        if code:
            # Patcher les appels .save() pour utiliser output_path
            # Le LLM écrit souvent doc.save("fichier.docx") au lieu de doc.save(output_path)
            code = re.sub(
                r'\.save\s*\(\s*["\'].*?["\']\s*\)',
                '.save(output_path)',
                code,
            )
            # Aussi patcher wb.save(...) pour xlsx
            code = re.sub(
                r'\.save\s*\(\s*(?:f["\'].*?["\']|["\'].*?["\'])\s*\)',
                '.save(output_path)',
                code,
            )

            # Ajouter .save(output_path) si absent (code tronqué par max_tokens
            # mais syntaxiquement valide, donc pas réparé par repair_truncated_code)
            code = _ensure_save_call(code)

            try:
                logger.info(
                    f"[{self.skill_id}] Code Python détecté, exécution sandboxée..."
                )
                await execute_sandboxed(
                    code=code,
                    output_path=str(output_path),
                    title=params.title,
                    format_type=self.output_format.value,
                    nb_slides=params.metadata.get("nb_slides", 10),
                )

                # Vérifier que le fichier a été créé et contient du contenu
                if output_path.exists() and output_path.stat().st_size > 0:
                    # BUG-043 : valider que le document n'est pas quasi vide
                    if not _validate_document_content(
                        str(output_path), self.output_format.value
                    ):
                        logger.warning(
                            f"[{self.skill_id}] Code exécuté mais document quasi vide "
                            f"(contenu insuffisant), fallback vers parser legacy"
                        )
                    else:
                        file_size = output_path.stat().st_size
                        logger.info(
                            f"[{self.skill_id}] Code-execution réussi : "
                            f"{output_path} ({file_size} bytes)"
                        )
                        return SkillResult(
                            file_id=file_id,
                            file_path=output_path,
                            file_name=output_path.name,
                            file_size=file_size,
                            mime_type=self.get_mime_type(),
                            format=self.output_format,
                        )
                else:
                    logger.warning(
                        f"[{self.skill_id}] Code exécuté mais fichier vide ou inexistant, "
                        f"fallback vers parser legacy"
                    )
            except CodeExecutionError as e:
                logger.warning(
                    f"[{self.skill_id}] Échec code-execution : {e}, "
                    f"fallback vers parser legacy"
                )
            except Exception as e:
                logger.warning(
                    f"[{self.skill_id}] Erreur inattendue code-execution : {e}, "
                    f"fallback vers parser legacy"
                )
        else:
            logger.info(
                f"[{self.skill_id}] Pas de bloc Python détecté, "
                f"fallback vers parser legacy"
            )

        # Nettoyer les blocs de code du contenu avant le fallback
        # pour éviter que le parser Markdown ne les rende comme du texte brut
        # On ne supprime que les blocs complets (avec ``` fermant)
        # pour ne pas perdre le contenu quand le bloc est tronqué
        cleaned_content = re.sub(
            r"```(?:python|py|javascript|js|bash|json)?\s*\n.*?```",
            "",
            params.content,
            flags=re.DOTALL,
        ).strip()
        # Si le contenu nettoyé est vide (tout était dans un bloc code),
        # on garde le contenu original - le parser docx_generator
        # filtrera les lignes de code via sa propre détection
        if cleaned_content and len(cleaned_content) > 50:
            params = SkillParams(
                title=params.title,
                content=cleaned_content,
                template=params.template,
                metadata=params.metadata,
            )

        # Fallback vers l'ancien parser
        return await self._fallback_execute(params, file_id, output_path)

    def get_markdown_prompt_addition(self) -> str:
        """
        Instructions alternatives pour les modèles incapables de générer du code Python.

        Demande du Markdown structuré au lieu de code python-docx/pptx/openpyxl.
        Peut être surchargé par chaque generator pour des instructions spécifiques.
        """
        return """
Génère le contenu en Markdown bien structuré.
Utilise : # Titre, ## Sections, ### Sous-sections, listes, **gras**, *italique*, tableaux Markdown.
NE génère PAS de code Python. Écris directement le contenu textuel.
"""

    @abstractmethod
    async def _fallback_execute(
        self, params: SkillParams, file_id: str, output_path: Path
    ) -> SkillResult:
        """
        Fallback vers l'ancien parser si le code-execution échoue.

        Args:
            params: Paramètres de génération
            file_id: ID du fichier pré-généré
            output_path: Chemin de sortie pré-calculé

        Returns:
            Résultat avec chemin vers le fichier généré
        """
        pass
