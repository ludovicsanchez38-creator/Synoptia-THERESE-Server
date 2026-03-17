"""
THERESE v2 - PowerPoint Generator Skill

Génère des présentations PowerPoint (.pptx) avec le style Synoptia.
Approche code-execution avec fallback parser legacy.
"""

import logging
import re
from pathlib import Path
from typing import Any

from app.services.skills.base import FileFormat, SkillParams, SkillResult
from app.services.skills.code_executor import CodeGenSkill
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)


# Palette Synoptia
SYNOPTIA_COLORS = {
    "background": RGBColor(0x0B, 0x12, 0x26),
    "surface": RGBColor(0x13, 0x1B, 0x35),
    "text": RGBColor(0xE6, 0xED, 0xF7),
    "muted": RGBColor(0xA9, 0xB8, 0xD8),
    "primary": RGBColor(0x17, 0x33, 0xA6),
    "accent_cyan": RGBColor(0x22, 0xD3, 0xEE),
    "accent_magenta": RGBColor(0xE1, 0x1D, 0x8D),
}


class PptxSkill(CodeGenSkill):
    """
    Skill de génération de présentations PowerPoint.

    Crée des présentations .pptx professionnelles avec le style Synoptia dark.
    Approche code-execution : le LLM génère du code python-pptx.
    Fallback automatique vers l'ancien parser Markdown.
    """

    skill_id = "pptx-pro"
    name = "Présentation PowerPoint"
    description = "Génère une présentation PowerPoint avec le style Synoptia"
    output_format = FileFormat.PPTX

    def __init__(self, output_dir: Path):
        super().__init__(output_dir)

    def get_system_prompt_addition(self) -> str:
        """Instructions pour le LLM : générer du code Python python-pptx."""
        return """
## Instructions pour génération de présentation PowerPoint

Tu dois générer un **bloc de code Python** complet utilisant la bibliothèque `python-pptx`.
Le code sera exécuté dans un environnement sandboxé avec les variables suivantes pré-injectées :
- `output_path` (str) : chemin où sauvegarder le fichier .pptx
- `title` (str) : sujet demandé par l'utilisateur (utilise-le comme inspiration, mais génère un titre professionnel et pertinent pour la slide de titre, PAS le prompt brut)
- `nb_slides` (int) : nombre de slides demandé par l'utilisateur. RESPECTE exactement ce nombre (slide de titre + slides de contenu + slide de fin = nb_slides total)
- `SYNOPTIA_COLORS` (dict) : palette de couleurs Synoptia

### Imports disponibles
pptx (Presentation), pptx.util (Inches, Pt, Cm, Emu), pptx.dml.color (RGBColor), pptx.enum.text (PP_ALIGN, MSO_ANCHOR), pptx.enum.shapes, datetime, json, re, math, Decimal

### Règles impératives
1. **Format** : 16:9 (13.333 x 7.5 inches)
   ```python
   prs.slide_width = Inches(13.333)
   prs.slide_height = Inches(7.5)
   ```
2. **Dark theme Synoptia** :
   - Fond : #0B1226 (background)
   - Texte principal : #E6EDF7
   - Accent cyan : #22D3EE (titres, barres, éléments clés)
   - Accent magenta : #E11D8D (chiffres importants)
   - Surface : #131B35 (cartes, encadrés)
3. **Layout** : Toujours utiliser blank layout (index 6), pas de placeholders pré-définis
   ```python
   blank_layout = prs.slide_layouts[6]
   slide = prs.slides.add_slide(blank_layout)
   ```
4. **Fond de chaque slide** :
   ```python
   background = slide.background
   fill = background.fill
   fill.solid()
   fill.fore_color.rgb = RGBColor(0x0B, 0x12, 0x26)
   ```
5. **Titre de slide** : Outfit, 36pt, bold, couleur RGBColor(0xE6, 0xED, 0xF7) (text_primary #E6EDF7), position top-left
6. **Barre accent cyan** sous chaque titre :
   ```python
   bar = slide.shapes.add_shape(1, Inches(0.75), Inches(1.4), Inches(2), Inches(0.05))
   bar.fill.solid()
   bar.fill.fore_color.rgb = RGBColor(0x22, 0xD3, 0xEE)
   bar.line.fill.background()
   ```
7. **Variété de slides** (ne pas faire que des bullet points) :
   - Slide titre : titre centré 54pt + sous-titre
   - Slide contenu : titre + 3-5 bullet points
   - Slide 2 colonnes : split gauche/droite pour comparaisons
   - Slide chiffres clés : 2-4 gros chiffres avec labels
   - Slide citation : texte centré en italique avec guillemets
   - Slide tableau : grille de données
8. **Texte** : Inter 24pt pour le corps, ne pas dépasser 6 lignes par slide
9. **Slide de fin** : "Merci" en 48pt cyan centré + "Généré par THÉRÈSE - Synoptïa" en muted
10. **Finir par** : `prs.save(output_path)`

**INTERDIT** : Ne jamais mettre de balises Markdown (`**`, `*`, `~~`, `_`) dans le texte des slides. Texte brut uniquement.

**INTERDIT** : Ne pas ajouter de bloc récapitulatif (Sujet / Action / Date) après la présentation. Le code se termine par `prs.save(output_path)` sans aucun autre output.

### Structure du code

```python
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper pour fond sombre
def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x0B, 0x12, 0x26)

# Slide 1 : Titre
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
set_bg(slide)
# ... titre centré ...

# Slides de contenu
# ... variété de layouts ...

# Slide finale
slide = prs.slides.add_slide(blank)
set_bg(slide)
# ... "Merci" + footer "Généré par THÉRÈSE - Synoptïa" ...

prs.save(output_path)
```

Génère UNIQUEMENT le bloc ```python``` avec le code complet. Pas d'explication avant ou après.
"""

    def get_markdown_prompt_addition(self) -> str:
        """Instructions Markdown pour modèles non code-capable."""
        return """
## Instructions pour génération de présentation PowerPoint

Génère le contenu de la présentation, slide par slide.
Sépare chaque slide par --- (trois tirets sur une ligne seule).
Chaque slide commence par ## Titre de la slide, suivi de 3-5 points clés.

### Structure attendue
1. Slide titre : ## Titre de la présentation (+ sous-titre en italique)
2. Slides de contenu : ## Titre + liste de points clés (- point)
3. Slide de fin : ## Merci

### Exemple de format
## Introduction au sujet
- Premier point important
- Deuxième point avec **élément clé**
- Troisième point

---

## Analyse détaillée
- Point d'analyse 1
- Point d'analyse 2

NE génère PAS de code Python. Écris directement le contenu textuel des slides.
"""

    async def _fallback_execute(
        self, params: SkillParams, file_id: str, output_path: Path
    ) -> SkillResult:
        """
        Fallback : ancien parser Markdown -> python-pptx.

        Args:
            params: Paramètres de génération
            file_id: ID du fichier pré-généré
            output_path: Chemin de sortie pré-calculé

        Returns:
            Résultat avec chemin vers le fichier généré
        """
        # Créer la présentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9
        prs.slide_height = Inches(7.5)

        # Parser le contenu et créer les slides
        slides_content = self._parse_content(params.content)

        # Slide de titre
        self._add_title_slide(prs, params.title)

        # Slides de contenu
        for slide_data in slides_content:
            self._add_content_slide(prs, slide_data)

        # Slide de fin
        self._add_end_slide(prs)

        # Sauvegarder
        prs.save(str(output_path))

        # Calculer la taille
        file_size = output_path.stat().st_size

        logger.info(f"Generated PPTX (fallback): {output_path} ({file_size} bytes)")

        return SkillResult(
            file_id=file_id,
            file_path=output_path,
            file_name=output_path.name,
            file_size=file_size,
            mime_type=self.get_mime_type(),
            format=self.output_format,
        )

    def _parse_content(self, content: str) -> list[dict[str, Any]]:
        """
        Parse le contenu en structure de slides.

        Args:
            content: Contenu généré par le LLM

        Returns:
            Liste de dictionnaires avec titre et points pour chaque slide
        """
        slides = []

        # Découper par délimiteurs --- ou par titres #
        blocks = re.split(r'\n---\n|\n-{3,}\n', content)

        for block in blocks:
            block = block.strip()
            if not block:
                continue

            lines = block.split('\n')
            title = None
            points = []

            for line in lines:
                line = line.strip()
                if not line:
                    continue

                # Titre de slide
                if line.startswith('# '):
                    title = line[2:].strip()
                elif line.startswith('## '):
                    title = line[3:].strip()
                # Points
                elif line.startswith('- ') or line.startswith('* '):
                    points.append(line[2:].strip())
                elif re.match(r'^\d+\.\s', line):
                    points.append(re.sub(r'^\d+\.\s', '', line).strip())
                # Texte normal (traiter comme point si pas de titre encore)
                elif title and line:
                    points.append(line)

            if title or points:
                slides.append({
                    "title": title or "Slide",
                    "points": points[:6],  # Max 6 points par slide
                })

        return slides

    def _add_title_slide(self, prs: Presentation, title: str) -> None:
        """Ajoute la slide de titre."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Fond sombre
        self._set_background(slide, SYNOPTIA_COLORS["background"])

        # Titre centré
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = title
        title_frame.paragraphs[0].font.name = "Outfit"
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["text"]
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Sous-titre
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.paragraphs[0].text = "Synoptia - L'entrepreneur augmenté"
        subtitle_frame.paragraphs[0].font.name = "Inter"
        subtitle_frame.paragraphs[0].font.size = Pt(18)
        subtitle_frame.paragraphs[0].font.italic = True
        subtitle_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["accent_cyan"]
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs: Presentation, slide_data: dict[str, Any]) -> None:
        """Ajoute une slide de contenu."""
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Fond sombre
        self._set_background(slide, SYNOPTIA_COLORS["background"])

        # Titre
        title_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(0.5), Inches(11.833), Inches(1)
        )
        title_frame = title_box.text_frame
        title_frame.paragraphs[0].text = slide_data["title"]
        title_frame.paragraphs[0].font.name = "Outfit"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["text"]

        # Barre accent sous le titre
        accent_bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.75), Inches(1.4), Inches(2), Inches(0.05)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = SYNOPTIA_COLORS["accent_cyan"]
        accent_bar.line.fill.background()

        # Points
        content_box = slide.shapes.add_textbox(
            Inches(0.75), Inches(1.8), Inches(11.833), Inches(5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for i, point in enumerate(slide_data.get("points", [])):
            if i == 0:
                p = content_frame.paragraphs[0]
            else:
                p = content_frame.add_paragraph()

            # Bullet avec accent
            p.text = f"  {point}"
            p.font.name = "Inter"
            p.font.size = Pt(24)
            p.font.color.rgb = SYNOPTIA_COLORS["text"]
            p.space_before = Pt(16)
            p.level = 0

    def _add_end_slide(self, prs: Presentation) -> None:
        """Ajoute la slide de fin."""
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Fond sombre
        self._set_background(slide, SYNOPTIA_COLORS["background"])

        # Texte de fin
        text_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(12.333), Inches(1.5)
        )
        text_frame = text_box.text_frame

        text_frame.paragraphs[0].text = "Merci"
        text_frame.paragraphs[0].font.name = "Outfit"
        text_frame.paragraphs[0].font.size = Pt(48)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = SYNOPTIA_COLORS["accent_cyan"]
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Sous-texte
        p = text_frame.add_paragraph()
        p.text = "Généré par THERESE - Synoptia"
        p.font.name = "Inter"
        p.font.size = Pt(14)
        p.font.color.rgb = SYNOPTIA_COLORS["muted"]
        p.alignment = PP_ALIGN.CENTER

    def _set_background(self, slide, color: RGBColor) -> None:
        """Définit la couleur de fond d'une slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color
